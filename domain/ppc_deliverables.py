"""
ppc_deliverables.py -- Polished, one-click PPC outputs.

Split out from ppc_module.py because these are FORMATTING code
(xlsx / docx / pptx / html), not PPC logic. Each function takes already-
computed data + config and produces a downloadable file.

Contents:
  1. build_status_xlsx      -- coloured Search Term Report status sheet
  2. build_audit_docx       -- PPC health audit (findings + prescriptions)
  3. build_dashboard_html   -- interactive HTML control-room dashboard
  4. build_forecast_xlsx    -- 3-scenario revenue + units forecast
  5. build_weekly_deck_pptx -- 5-slide weekly performance deck

HONEST SCOPE: these are functional MVPs, not the full skill outputs from
the handover doc. Each is one hand-built pass; the deep skill versions
(hundreds of lines each) get built as they become the top priority.
"""
from __future__ import annotations
import io, os, statistics
from typing import Optional, Any


# =============================================================================
# 1. STATUS XLSX -- coloured Search Term Report status sheet
# =============================================================================
def build_status_xlsx(status_rows: list, currency: str = "£") -> bytes:
    """Emit a coloured .xlsx status sheet: one row per search term with
    metrics + colour-coded status column (green/amber/red). Uses openpyxl
    with an actual PatternFill so the colours travel with the file."""
    from openpyxl import Workbook
    from openpyxl.styles import PatternFill, Font, Alignment
    from openpyxl.utils import get_column_letter

    wb = Workbook()
    ws = wb.active
    ws.title = "Search Term Status"

    cols = [
        ("customer_search_term", "Search Term",     28),
        ("triggering_keyword",   "Triggering KW",   20),
        ("match_type",           "Match Type",      10),
        ("campaign",             "Campaign",        28),
        ("ad_group",             "Ad Group",        22),
        ("impressions",          "Impr",             8),
        ("clicks",               "Clicks",           7),
        ("ctr",                  "CTR",              7),
        ("cpc",                  "CPC",              8),
        ("spend",                "Spend",            9),
        ("orders",               "Orders",           7),
        ("units",                "Units",            7),
        ("cvr",                  "CVR",              7),
        ("sales",                "Sales",            9),
        ("acos",                 "ACOS",             7),
        ("roas",                 "ROAS",             6),
        ("status",               "Status",          22),
    ]

    # Header row
    header_font = Font(bold=True, color="FFFFFF")
    header_fill = PatternFill("solid", fgColor="223B65")
    for i, (_, label, _) in enumerate(cols, 1):
        cell = ws.cell(row=1, column=i, value=label)
        cell.font = header_font
        cell.fill = header_fill
        cell.alignment = Alignment(vertical="center", horizontal="center")
        ws.column_dimensions[get_column_letter(i)].width = cols[i-1][2]
    ws.freeze_panes = "A2"

    # Status -> colour map. Muted to be readable on printouts too.
    STATUS_FILLS = {
        "CONVERTING":              PatternFill("solid", fgColor="C8E6C9"),  # green
        "CONVERTS-BUT-HIGH-ACOS":  PatternFill("solid", fgColor="FFE0B2"),  # amber
        "HIGH-SPEND-WATCH":        PatternFill("solid", fgColor="FFE0B2"),  # amber
        "OVER-$10-CUT":            PatternFill("solid", fgColor="FFCDD2"),  # red
        "CLICKS-NO-SALE":          PatternFill("solid", fgColor="FFCDD2"),  # red
        "EARLY":                   PatternFill("solid", fgColor="E3F2FD"),  # neutral blue
        "IMPRESSIONS-ONLY":        PatternFill("solid", fgColor="F5F5F5"),  # grey
    }
    default_fill = PatternFill("solid", fgColor="FFFFFF")

    # Data rows
    for r_idx, row in enumerate(status_rows, 2):
        fill = STATUS_FILLS.get(row.get("status", ""), default_fill)
        for c_idx, (key, _, _) in enumerate(cols, 1):
            v = row.get(key, "")
            # Type-aware formatting
            if key in ("ctr", "cvr", "acos") and isinstance(v, (int, float)):
                cell = ws.cell(row=r_idx, column=c_idx, value=v)
                cell.number_format = "0.0%"
            elif key in ("cpc", "spend", "sales") and isinstance(v, (int, float)):
                cell = ws.cell(row=r_idx, column=c_idx, value=v)
                # openpyxl's currency format uses locale, so just show it as number w/ 2dp
                cell.number_format = f'"{currency}"#,##0.00'
            elif key == "roas" and isinstance(v, (int, float)):
                cell = ws.cell(row=r_idx, column=c_idx, value=v)
                cell.number_format = "0.00"
            elif key in ("impressions", "clicks", "orders", "units") and isinstance(v, (int, float)):
                cell = ws.cell(row=r_idx, column=c_idx, value=int(v))
                cell.number_format = "#,##0"
            else:
                cell = ws.cell(row=r_idx, column=c_idx, value=v)
            # Colour the whole row lightly; status column gets a stronger colour
            if key == "status":
                cell.fill = fill
                cell.font = Font(bold=True)
                cell.alignment = Alignment(horizontal="center")
            else:
                # subtle tint: use a lighter version by only filling the status cell strongly
                pass

    out = io.BytesIO()
    wb.save(out)
    return out.getvalue()


# =============================================================================
# 2. AUDIT DOCX -- PPC health audit
# =============================================================================
def _audit_findings_from_bulk(bulk_rows: list) -> dict:
    """Compute structural findings from an SP bulk export. Returns a dict of
    findings the docx builder turns into sections. Every finding is either a
    hard rule violation (from the handover doc) or a numeric benchmark check."""
    findings = {
        "campaign_count":        0,
        "adgroup_count":         0,
        "keyword_count":         0,
        "product_target_count":  0,
        "total_daily_budget":    0.0,
        "match_type_pyramid":    {"Exact": 0, "Phrase": 0, "Broad": 0},
        "kw_pt_shared_adgroups": [],
        "duplicate_kw_pairs":    [],
        "campaigns_without_ads": [],
        "campaigns_missing_match_types": [],
    }
    campaigns_by_name    = {}
    adgroups_kw_matches  = {}   # (camp, ag) -> set of match types
    kw_ag_set            = set()
    pt_ag_set            = set()
    seen_kw_pairs        = set()
    campaigns_with_ads   = set()
    all_campaigns        = set()

    for r in bulk_rows:
        entity = r.get("Entity") or r.get("entity") or ""
        camp   = r.get("Campaign Name") or r.get("campaign") or ""
        ag     = r.get("Ad Group Name") or r.get("ad_group") or ""

        if entity == "Campaign":
            findings["campaign_count"] += 1
            all_campaigns.add(camp)
            budget = r.get("Daily Budget") or r.get("daily_budget") or 0
            try:
                findings["total_daily_budget"] += float(budget)
            except (TypeError, ValueError):
                pass
        elif entity == "Ad Group":
            findings["adgroup_count"] += 1
        elif entity == "Keyword":
            findings["keyword_count"] += 1
            mt = r.get("Match Type") or r.get("match_type") or ""
            if mt in findings["match_type_pyramid"]:
                findings["match_type_pyramid"][mt] += 1
            k = (r.get("Keyword Text", "").lower(), mt.lower())
            if k in seen_kw_pairs:
                findings["duplicate_kw_pairs"].append(k)
            seen_kw_pairs.add(k)
            kw_ag_set.add((camp, ag))
            adgroups_kw_matches.setdefault((camp, ag), set()).add(mt)
        elif entity == "Product Targeting":
            findings["product_target_count"] += 1
            pt_ag_set.add((camp, ag))
        elif entity == "Product Ad":
            campaigns_with_ads.add(camp)

    findings["kw_pt_shared_adgroups"] = sorted(kw_ag_set & pt_ag_set)
    findings["campaigns_without_ads"] = sorted(all_campaigns - campaigns_with_ads)
    # match-type coverage per ad group (base rule: all 3 present portfolio-wide,
    # not per ad group -- so we surface ad groups running only 1 match type as
    # informational, not as a violation)
    findings["adgroups_single_match_type"] = [
        (c, a, sorted(ms)) for (c, a), ms in adgroups_kw_matches.items() if len(ms) == 1
    ]
    return findings


def build_audit_docx(bulk_rows: list,
                     performance_rows: Optional[list] = None,
                     account_label: str = "",
                     marketplace: str = "UK") -> bytes:
    """Produce a PPC health-audit .docx from the SP bulk export.
    Optional performance_rows adds a spend/sales section; without it the audit
    is structural-only (no financial claims made from missing data).
    """
    from docx import Document
    from docx.shared import Inches, Pt, RGBColor

    doc = Document()

    # Title
    t = doc.add_heading(f"Amazon PPC Audit -- {account_label or 'Account'}", level=0)
    doc.add_paragraph(f"Marketplace: {marketplace}   |   Scope: Sponsored Products bulk export")

    findings = _audit_findings_from_bulk(bulk_rows)

    # Executive summary
    doc.add_heading("Executive summary", level=1)
    doc.add_paragraph(
        f"Portfolio has {findings['campaign_count']} campaign(s), "
        f"{findings['adgroup_count']} ad group(s), "
        f"{findings['keyword_count']} keyword(s), "
        f"{findings['product_target_count']} product target(s). "
        f"Total configured daily budget: £{findings['total_daily_budget']:.2f}."
    )

    # HARD-RULE violations (from the doc's non-negotiables)
    doc.add_heading("Hard-rule checks", level=1)

    # 1. keywords + product targets in same ad group
    p = doc.add_paragraph()
    if findings["kw_pt_shared_adgroups"]:
        r = p.add_run("VIOLATION: Keywords and Product Targets share ad groups. ")
        r.bold = True; r.font.color.rgb = RGBColor(0xC0, 0x21, 0x21)
        p.add_run("The handover doc requires these in SEPARATE ad groups because "
                  "Amazon bids independently by targeting type and mixing them "
                  "corrupts your bid signals. Affected:")
        for (c, a) in findings["kw_pt_shared_adgroups"]:
            doc.add_paragraph(f"{c} / {a}", style="List Bullet")
    else:
        r = p.add_run("OK: no ad group mixes keywords with product targets.")
        r.font.color.rgb = RGBColor(0x1B, 0x5E, 0x20)

    # 2. duplicate (keyword, match-type) pairs
    p = doc.add_paragraph()
    if findings["duplicate_kw_pairs"]:
        r = p.add_run(f"VIOLATION: {len(findings['duplicate_kw_pairs'])} duplicate "
                      f"(keyword, match-type) pair(s). ")
        r.bold = True; r.font.color.rgb = RGBColor(0xC0, 0x21, 0x21)
        p.add_run("Amazon lets you upload duplicates but they cannibalise your "
                  "own bids at auction, quietly inflating CPC.")
    else:
        r = p.add_run("OK: no duplicate keyword/match-type pairs across the portfolio.")
        r.font.color.rgb = RGBColor(0x1B, 0x5E, 0x20)

    # 3. campaigns without product ads (dead campaigns)
    p = doc.add_paragraph()
    if findings["campaigns_without_ads"]:
        r = p.add_run(f"WARNING: {len(findings['campaigns_without_ads'])} campaign(s) "
                      f"have no Product Ad row. ")
        r.bold = True; r.font.color.rgb = RGBColor(0xE6, 0x5C, 0x00)
        p.add_run("Without a Product Ad the campaign can't serve. Likely intentional "
                  "(paused/orphaned) but worth confirming:")
        for c in findings["campaigns_without_ads"]:
            doc.add_paragraph(c, style="List Bullet")
    else:
        r = p.add_run("OK: every campaign has at least one Product Ad.")
        r.font.color.rgb = RGBColor(0x1B, 0x5E, 0x20)

    # Structural benchmarks
    doc.add_heading("Match-type pyramid", level=1)
    p = findings["match_type_pyramid"]
    total = sum(p.values()) or 1
    doc.add_paragraph(
        f"Exact: {p['Exact']} ({p['Exact']*100/total:.0f}%) | "
        f"Phrase: {p['Phrase']} ({p['Phrase']*100/total:.0f}%) | "
        f"Broad: {p['Broad']} ({p['Broad']*100/total:.0f}%)"
    )
    doc.add_paragraph(
        "Healthy full-coverage portfolios keep the three types roughly balanced "
        "across the account (each keyword ideally exists in all three). If Exact "
        "is over 60% you're likely under-discovering new terms; if Broad is over "
        "50% you're likely burning waste."
    )

    # Ad groups running only 1 match type (informational)
    if findings["adgroups_single_match_type"]:
        doc.add_heading("Ad groups running only one match type", level=2)
        doc.add_paragraph(
            f"{len(findings['adgroups_single_match_type'])} ad group(s) contain a "
            "single match type. This is fine when it's intentional (e.g. an "
            "Exact_GROW ad group holds only Exact), but flag them for review:"
        )
        for c, a, ms in findings["adgroups_single_match_type"][:20]:
            doc.add_paragraph(f"{c} / {a}: {ms}", style="List Bullet")
        if len(findings["adgroups_single_match_type"]) > 20:
            doc.add_paragraph(f"... and {len(findings['adgroups_single_match_type'])-20} more")

    # Performance section (only if performance data was supplied -- no fabrication)
    if performance_rows:
        doc.add_heading("Financial checks", level=1)
        spend = sum(_num(r.get("spend"))  for r in performance_rows)
        sales = sum(_num(r.get("sales")) for r in performance_rows)
        orders = sum(_num(r.get("orders")) for r in performance_rows)
        acos = round(spend/sales, 4) if sales > 0 else 0
        roas = round(sales/spend, 4) if spend > 0 else 0
        doc.add_paragraph(
            f"Total spend: {spend:.2f}   Sales: {sales:.2f}   "
            f"Orders: {int(orders)}   ACOS: {acos*100:.1f}%   ROAS: {roas:.2f}"
        )
        # Waste = zero-order campaigns
        by_camp = {}
        for r in performance_rows:
            c = r.get("campaign") or r.get("Campaign Name") or ""
            by_camp.setdefault(c, {"spend": 0, "orders": 0})
            by_camp[c]["spend"]  += _num(r.get("spend"))
            by_camp[c]["orders"] += _num(r.get("orders"))
        zero_order_spend = sum(v["spend"] for v in by_camp.values() if v["orders"] == 0)
        if zero_order_spend > 0:
            p = doc.add_paragraph()
            r = p.add_run(f"WARNING: {zero_order_spend:.2f} spent on campaigns with zero orders. ")
            r.bold = True; r.font.color.rgb = RGBColor(0xE6, 0x5C, 0x00)
            p.add_run("These are candidates for pause or restructure.")
    else:
        doc.add_heading("Financial checks", level=1)
        doc.add_paragraph(
            "No performance data was provided, so no financial claims are made. "
            "For a financial audit, upload a Sponsored Products campaign performance "
            "report (last 30 days) alongside the bulk export."
        )

    doc.add_heading("Recommended next actions", level=1)
    doc.add_paragraph("1. Fix any hard-rule violations flagged above before anything else.", style="List Number")
    doc.add_paragraph("2. Run a search-term harvest to promote proven converters + negate $10+ zero-order terms.", style="List Number")
    doc.add_paragraph("3. Rebalance match-type pyramid if it's skewed (Exact-heavy = add discovery; Broad-heavy = tighten with phrase/exact).", style="List Number")

    out = io.BytesIO()
    doc.save(out)
    return out.getvalue()


# =============================================================================
# 3. DASHBOARD HTML -- self-contained interactive control-room
# =============================================================================
def build_dashboard_html(bulk_rows: list,
                          performance_rows: Optional[list] = None,
                          account_label: str = "",
                          marketplace: str = "UK") -> bytes:
    """Emit a single-file interactive dashboard as HTML with inline CSS + JS.
    Panels: overview, spend distribution, match-type pyramid, action list."""
    findings = _audit_findings_from_bulk(bulk_rows)

    # Aggregate performance if provided
    perf_by_camp = {}
    total_spend = total_sales = total_orders = 0.0
    if performance_rows:
        for r in performance_rows:
            c = r.get("campaign") or r.get("Campaign Name") or ""
            perf_by_camp.setdefault(c, {"spend": 0.0, "sales": 0.0, "orders": 0})
            s = _num(r.get("spend")); sa = _num(r.get("sales")); o = _num(r.get("orders"))
            perf_by_camp[c]["spend"]  += s
            perf_by_camp[c]["sales"]  += sa
            perf_by_camp[c]["orders"] += o
            total_spend  += s
            total_sales  += sa
            total_orders += o

    top_spend = sorted(perf_by_camp.items(), key=lambda kv: kv[1]["spend"], reverse=True)[:10]
    acos = round(total_spend/total_sales, 4) if total_sales > 0 else 0
    roas = round(total_sales/total_spend, 4) if total_spend > 0 else 0

    # Match-type pyramid data
    p = findings["match_type_pyramid"]
    ptotal = sum(p.values()) or 1

    # Action items
    actions = []
    if findings["kw_pt_shared_adgroups"]:
        actions.append(("HARD RULE VIOLATION", "red",
                        f"{len(findings['kw_pt_shared_adgroups'])} ad group(s) mix keywords with product targets. Split them."))
    if findings["duplicate_kw_pairs"]:
        actions.append(("HARD RULE VIOLATION", "red",
                        f"{len(findings['duplicate_kw_pairs'])} duplicate (keyword, match-type) pair(s). Remove duplicates."))
    if findings["campaigns_without_ads"]:
        actions.append(("Warning", "amber",
                        f"{len(findings['campaigns_without_ads'])} campaign(s) have no Product Ad and can't serve."))
    zero_ord_spend = sum(v["spend"] for v in perf_by_camp.values() if v["orders"] == 0) if perf_by_camp else 0
    if zero_ord_spend > 0:
        actions.append(("Waste", "amber",
                        f"£{zero_ord_spend:.2f} spent on campaigns with zero orders. Candidates for pause."))
    if not actions:
        actions.append(("OK", "green", "No hard-rule violations or obvious waste detected."))

    html = f"""<!DOCTYPE html>
<html lang="en"><head><meta charset="utf-8">
<title>PPC Control Room — {_esc(account_label)}</title>
<style>
  body{{margin:0;padding:24px;font-family:-apple-system,BlinkMacSystemFont,Segoe UI,sans-serif;background:#0f1420;color:#e8eaed}}
  h1{{margin:0 0 4px;font-size:22px;font-weight:700}}
  .sub{{opacity:.7;font-size:13px;margin-bottom:24px}}
  .grid{{display:grid;grid-template-columns:repeat(auto-fit,minmax(280px,1fr));gap:16px}}
  .panel{{background:#141b2b;border:1px solid #263145;border-radius:10px;padding:16px}}
  .panel h2{{margin:0 0 12px;font-size:14px;font-weight:600;opacity:.85;text-transform:uppercase;letter-spacing:.5px}}
  .kpi{{display:flex;gap:24px;flex-wrap:wrap;margin-bottom:8px}}
  .kpi div{{min-width:80px}}
  .kpi .lbl{{font-size:11px;opacity:.65;text-transform:uppercase;letter-spacing:.5px}}
  .kpi .val{{font-size:20px;font-weight:700}}
  table{{width:100%;border-collapse:collapse;font-size:13px}}
  th,td{{text-align:left;padding:6px 4px;border-bottom:1px solid #263145}}
  th{{opacity:.7;font-weight:600;font-size:11px;text-transform:uppercase}}
  .bar{{height:16px;border-radius:4px;background:#243154;position:relative;margin:4px 0}}
  .bar .fill{{position:absolute;top:0;bottom:0;left:0;border-radius:4px;background:#3a7cff}}
  .bar .fill.exact{{background:#4caf50}}
  .bar .fill.phrase{{background:#ffb74d}}
  .bar .fill.broad{{background:#e57373}}
  .action{{padding:10px;border-radius:8px;margin:6px 0;font-size:13px}}
  .action.red{{background:#3a1f1f;border:1px solid #7a2a2a;color:#ff8a8a}}
  .action.amber{{background:#3a2f1a;border:1px solid #7a5a2a;color:#ffce7a}}
  .action.green{{background:#1c3a1c;border:1px solid #2a7a2a;color:#8adca0}}
  .lbl-inline{{opacity:.65;font-size:11px;margin-right:8px}}
</style></head><body>
<h1>PPC Control Room — {_esc(account_label or 'Account')}</h1>
<div class="sub">Marketplace: {_esc(marketplace)} · Portfolio scope from bulk export</div>

<div class="grid">
  <div class="panel">
    <h2>Portfolio</h2>
    <div class="kpi">
      <div><div class="lbl">Campaigns</div><div class="val">{findings['campaign_count']}</div></div>
      <div><div class="lbl">Ad groups</div><div class="val">{findings['adgroup_count']}</div></div>
      <div><div class="lbl">Keywords</div><div class="val">{findings['keyword_count']}</div></div>
      <div><div class="lbl">Product tgts</div><div class="val">{findings['product_target_count']}</div></div>
    </div>
    <div style="margin-top:8px"><span class="lbl-inline">Daily budget:</span><b>£{findings['total_daily_budget']:.2f}</b></div>
  </div>

  <div class="panel">
    <h2>Performance {'(from performance report)' if performance_rows else '(no data)'}</h2>
    {'''<div class="kpi">
      <div><div class="lbl">Spend</div><div class="val">£{spend:.2f}</div></div>
      <div><div class="lbl">Sales</div><div class="val">£{sales:.2f}</div></div>
      <div><div class="lbl">Orders</div><div class="val">{orders}</div></div>
      <div><div class="lbl">ACOS</div><div class="val">{acos:.1f}%</div></div>
      <div><div class="lbl">ROAS</div><div class="val">{roas:.2f}</div></div>
    </div>'''.format(spend=total_spend, sales=total_sales, orders=int(total_orders),
                     acos=acos*100, roas=roas) if performance_rows else
     '<div style="opacity:.6;font-size:13px">Upload a campaign performance report to unlock this panel.</div>'}
  </div>

  <div class="panel">
    <h2>Match-type pyramid</h2>
    <div><div class="lbl-inline">Exact</div><b>{p['Exact']}</b> ({p['Exact']*100/ptotal:.0f}%)
      <div class="bar"><div class="fill exact" style="width:{p['Exact']*100/ptotal}%"></div></div>
    </div>
    <div><div class="lbl-inline">Phrase</div><b>{p['Phrase']}</b> ({p['Phrase']*100/ptotal:.0f}%)
      <div class="bar"><div class="fill phrase" style="width:{p['Phrase']*100/ptotal}%"></div></div>
    </div>
    <div><div class="lbl-inline">Broad</div><b>{p['Broad']}</b> ({p['Broad']*100/ptotal:.0f}%)
      <div class="bar"><div class="fill broad" style="width:{p['Broad']*100/ptotal}%"></div></div>
    </div>
  </div>

  <div class="panel" style="grid-column:1/-1">
    <h2>Action items</h2>
    {"".join(f'<div class="action {col}"><b>[{lbl}]</b> {_esc(msg)}</div>' for lbl,col,msg in actions)}
  </div>

  {('<div class="panel" style="grid-column:1/-1"><h2>Top 10 campaigns by spend</h2><table>'
     '<tr><th>Campaign</th><th>Spend</th><th>Sales</th><th>ACOS</th><th>Orders</th></tr>'
     + "".join(f'<tr><td>{_esc(c)}</td><td>£{v["spend"]:.2f}</td><td>£{v["sales"]:.2f}</td>'
                f'<td>{(v["spend"]/v["sales"]*100 if v["sales"]>0 else 0):.1f}%</td>'
                f'<td>{int(v["orders"])}</td></tr>' for c,v in top_spend)
     + '</table></div>') if top_spend else ''}
</div>
</body></html>"""
    return html.encode("utf-8")


# =============================================================================
# 4. FORECAST XLSX -- 3-scenario revenue + units
# =============================================================================
def build_forecast_xlsx(business_report_rows: list,
                        target_tacos_pct: float = 15.0,
                        net_margin_pct:   float = 25.0,
                        months_ahead: int = 3) -> bytes:
    """Build a 3-scenario forecast xlsx from an Amazon Business Report.

    Uses ORDERED-UNITS (last N days) as the base, projects forward using three
    scenarios: baseline (flat), conservative (-10% then decay), aggressive (+15%).
    Applies target TACOS to compute suggested ad spend and expected profit at
    the given net margin.

    HONEST: this is a simple MVP. Advanced forecast (weekly seasonality, event
    detection, elasticity curves) belongs in the amazon-sales-projections skill.
    """
    from openpyxl import Workbook
    from openpyxl.styles import PatternFill, Font, Alignment

    # Extract ordered_product_sales + units. Business report column names vary by year.
    ORDERS_FIELDS = ("ordered_product_sales", "ordered product sales", "sales", "ordered_product_sales_b2c")
    UNITS_FIELDS  = ("units_ordered", "units ordered", "units")

    total_sales = 0.0
    total_units = 0
    for r in business_report_rows:
        for f in ORDERS_FIELDS:
            v = r.get(f)
            if v is not None and v != "":
                total_sales += _num(v)
                break
        for f in UNITS_FIELDS:
            v = r.get(f)
            if v is not None and v != "":
                total_units += int(_num(v))
                break

    days_in_report = max(len(business_report_rows), 1)
    daily_sales = total_sales / days_in_report
    daily_units = total_units / days_in_report

    wb = Workbook()
    ws = wb.active
    ws.title = "3-scenario forecast"

    # Header block
    header_font = Font(bold=True, color="FFFFFF")
    header_fill = PatternFill("solid", fgColor="223B65")
    for col, label in enumerate(["Month", "Scenario", "Est. Units", "Est. Sales", "Suggested Ad Spend", "Est. Profit"], 1):
        c = ws.cell(row=1, column=col, value=label)
        c.font = header_font; c.fill = header_fill
        c.alignment = Alignment(horizontal="center")

    scenarios = [
        ("Baseline",     1.00),
        ("Conservative", 0.90),
        ("Aggressive",   1.15),
    ]

    row = 2
    days_per_month = 30
    for m in range(1, months_ahead + 1):
        for scenario_name, mult in scenarios:
            u = round(daily_units * days_per_month * mult ** m)
            s = round(daily_sales * days_per_month * mult ** m, 2)
            ad = round(s * target_tacos_pct / 100, 2)
            profit = round(s * net_margin_pct / 100 - ad, 2)

            ws.cell(row=row, column=1, value=f"Month {m}")
            ws.cell(row=row, column=2, value=scenario_name)
            ws.cell(row=row, column=3, value=u)
            ws.cell(row=row, column=4, value=s); ws.cell(row=row, column=4).number_format = '"£"#,##0.00'
            ws.cell(row=row, column=5, value=ad); ws.cell(row=row, column=5).number_format = '"£"#,##0.00'
            p = ws.cell(row=row, column=6, value=profit); p.number_format = '"£"#,##0.00'
            if profit < 0:
                p.fill = PatternFill("solid", fgColor="FFCDD2")
            row += 1
        row += 1        # gap between months

    # Notes tab -- explains assumptions honestly
    notes = wb.create_sheet("Assumptions")
    notes["A1"] = "Assumptions and honest caveats"
    notes["A1"].font = Font(bold=True, size=14)
    notes["A3"] = f"Report days used: {days_in_report}"
    notes["A4"] = f"Daily avg sales: £{daily_sales:.2f}"
    notes["A5"] = f"Daily avg units: {daily_units:.1f}"
    notes["A6"] = f"Target TACOS: {target_tacos_pct}%"
    notes["A7"] = f"Net margin: {net_margin_pct}%"
    notes["A9"] = "IMPORTANT: this MVP does not model weekly seasonality, Prime Day, "
    notes["A10"] = "BFCM, launch decay, or inventory constraints. Compound scaling per month is"
    notes["A11"] = "simple; long horizons will drift. For a rigorous forecast, use the"
    notes["A12"] = "amazon-sales-projections skill outside the app."

    out = io.BytesIO()
    wb.save(out)
    return out.getvalue()


# =============================================================================
# 5. WEEKLY-DECK PPTX -- 5-slide MVP
# =============================================================================
def build_weekly_deck_pptx(this_week_rows: list,
                            last_week_rows: Optional[list] = None,
                            brand: str = "Brand",
                            week_ending: str = "") -> bytes:
    """Build a 5-slide PPTX summarising this week vs last week's performance.
    Uses a simple layout so it opens cleanly everywhere without needing a
    template. Intentionally small; the full-fidelity deck is the weekly-client-deck skill.
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt
    # python-pptx keeps RGBColor under pptx.dml.color -- there's no pptx.dgm.color.
    # We import defensively so tests can still run in stripped-down envs.
    try:
        from pptx.dml.color import RGBColor as _RGBColor
    except ImportError:
        _RGBColor = None

    prs = Presentation()
    prs.slide_width  = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    def _kpis(rows):
        s = sum(_num(r.get("spend"))  for r in rows) if rows else 0
        sa = sum(_num(r.get("sales")) for r in rows) if rows else 0
        o = sum(_num(r.get("orders")) for r in rows) if rows else 0
        acos = (s/sa*100) if sa > 0 else 0
        roas = (sa/s)     if s  > 0 else 0
        return {"spend": s, "sales": sa, "orders": o, "acos": acos, "roas": roas}

    tw = _kpis(this_week_rows)
    lw = _kpis(last_week_rows) if last_week_rows else None

    def _add_title_slide():
        s = prs.slides.add_slide(blank)
        box = s.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(2))
        tf = box.text_frame
        tf.text = f"{brand} — Weekly PPC Review"
        tf.paragraphs[0].font.size = Pt(44)
        tf.paragraphs[0].font.bold = True
        p2 = tf.add_paragraph()
        p2.text = f"Week ending: {week_ending or '(unspecified)'}"
        p2.font.size = Pt(20)

    def _add_kpi_slide():
        s = prs.slides.add_slide(blank)
        title = s.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.7))
        title.text_frame.text = "This week — key numbers"
        title.text_frame.paragraphs[0].font.size = Pt(32); title.text_frame.paragraphs[0].font.bold = True

        labels = [("Spend", f"£{tw['spend']:.2f}"),
                  ("Sales", f"£{tw['sales']:.2f}"),
                  ("Orders", f"{int(tw['orders'])}"),
                  ("ACOS", f"{tw['acos']:.1f}%"),
                  ("ROAS", f"{tw['roas']:.2f}")]
        for i, (lbl, val) in enumerate(labels):
            x = Inches(0.5 + i * 2.5); y = Inches(2)
            box = s.shapes.add_textbox(x, y, Inches(2.3), Inches(2))
            tf = box.text_frame
            tf.text = lbl
            tf.paragraphs[0].font.size = Pt(16); tf.paragraphs[0].font.bold = False
            p2 = tf.add_paragraph(); p2.text = val
            p2.font.size = Pt(28); p2.font.bold = True

    def _add_wow_slide():
        s = prs.slides.add_slide(blank)
        title = s.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.7))
        title.text_frame.text = "Week-over-week"
        title.text_frame.paragraphs[0].font.size = Pt(32); title.text_frame.paragraphs[0].font.bold = True
        box = s.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(5))
        tf = box.text_frame
        if not lw:
            tf.text = "No prior-week data provided. Upload last week's performance to enable WoW comparison."
            tf.paragraphs[0].font.size = Pt(16)
            return
        def _delta(a, b, is_pct=False):
            if b == 0: return "n/a"
            d = (a - b) / b * 100
            arrow = "▲" if d > 0 else ("▼" if d < 0 else "▬")
            return f"{arrow} {abs(d):.1f}%"
        rows = [
            ("Spend",  f"£{lw['spend']:.2f}",  f"£{tw['spend']:.2f}",  _delta(tw['spend'],  lw['spend'])),
            ("Sales",  f"£{lw['sales']:.2f}",  f"£{tw['sales']:.2f}",  _delta(tw['sales'],  lw['sales'])),
            ("Orders", f"{int(lw['orders'])}", f"{int(tw['orders'])}", _delta(tw['orders'], lw['orders'])),
            ("ACOS",   f"{lw['acos']:.1f}%",   f"{tw['acos']:.1f}%",   f"{tw['acos']-lw['acos']:+.1f}pp"),
        ]
        tf.text = "Metric".ljust(10) + "Last".ljust(15) + "This".ljust(15) + "Δ"
        tf.paragraphs[0].font.size = Pt(14); tf.paragraphs[0].font.bold = True; tf.paragraphs[0].font.name = "Consolas"
        for m, a, b, d in rows:
            p = tf.add_paragraph(); p.text = m.ljust(10) + a.ljust(15) + b.ljust(15) + d
            p.font.size = Pt(14); p.font.name = "Consolas"

    def _add_next_slide():
        s = prs.slides.add_slide(blank)
        title = s.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.7))
        title.text_frame.text = "Next week"
        title.text_frame.paragraphs[0].font.size = Pt(32); title.text_frame.paragraphs[0].font.bold = True
        box = s.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(5.5))
        tf = box.text_frame
        tf.text = "Priorities (edit before sharing):"
        tf.paragraphs[0].font.size = Pt(18); tf.paragraphs[0].font.bold = True
        for line in ("1. Run search-term harvest and apply harvest bulk + negatives",
                     "2. Review any newly-flagged HIGH-ACOS terms; bid-down candidates for user approval",
                     "3. Verify no new hard-rule violations (kw + PT in same ad group; duplicate pairs)",
                     "4. Discuss any budget changes -- no auto-adjustments made"):
            p = tf.add_paragraph(); p.text = line; p.font.size = Pt(16)

    def _add_close_slide():
        s = prs.slides.add_slide(blank)
        box = s.shapes.add_textbox(Inches(2), Inches(3), Inches(9), Inches(1.5))
        tf = box.text_frame
        tf.text = "Questions?"
        tf.paragraphs[0].font.size = Pt(54); tf.paragraphs[0].font.bold = True

    _add_title_slide()
    _add_kpi_slide()
    _add_wow_slide()
    _add_next_slide()
    _add_close_slide()

    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()


# =============================================================================
# helpers
# =============================================================================
def _num(x, default=0.0) -> float:
    if x is None or x == "":
        return default
    if isinstance(x, (int, float)):
        return float(x)
    import re
    s = re.sub(r"[^\d\.\-]", "", str(x))
    try:
        return float(s) if s else default
    except ValueError:
        return default


def _esc(s: Any) -> str:
    return (str(s or "")
            .replace("&", "&amp;")
            .replace("<", "&lt;")
            .replace(">", "&gt;")
            .replace('"', "&quot;"))
